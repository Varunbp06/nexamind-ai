import hashlib
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from loguru import logger
from pairag.file.readers.base import BaseReader, FileItem, Document, List
from pairag.file.store.base import BaseFileStore
from pairag.file.utils.image_utils import compress_image_if_needed
from pairag.file.utils.markdown_tree_utils import PaiTable, convert_table_to_markdown
from pairag.file.utils.image_caption_tool import ImageCaptionTool
from pairag.file.utils.image_utils import to_markdown_image_text
from pptx.enum.shapes import PP_PLACEHOLDER
from pairag.file.utils.text_utils import replace_consecutive_spaces
import re


class PptxReader(BaseReader):
    def __init__(
        self, file_store: BaseFileStore, image_caption_tool: ImageCaptionTool = None
    ):
        self.file_store = file_store
        self.image_caption_tool = image_caption_tool
        logger.info("PptxReader inited.")


    def _extract_image_from_shape(self,shape, save_name_template: str, tenant_id: str):
        """从 shape 提取图片（支持 PICTURE、PLACEHOLDER.PICTURE、CHART 等）"""
        markdown = []
        images = []
        if not (hasattr(shape, 'image') and shape.image):
            return markdown, images
        if self.image_caption_tool:
            image_blob = shape.image.blob
            image_name = hashlib.md5(image_blob).hexdigest() + ".jpeg"
            save_image_name = save_name_template.format(image_name)
            image_file = BytesIO(image_blob)
            image_file = compress_image_if_needed(image_file)
            if image_file:
                try:
                    upload_result = self.file_store.write(file=image_file, file_name=image_name, file_path=save_image_name, tenant_id=tenant_id)
                    image_file.seek(0)
                    image_alt_text = self.image_caption_tool.extract_image(image_file.read())
                    if image_alt_text:
                        cleaned_alt = re.sub(r'\n', ' ', image_alt_text).replace('\r', '').strip()
                        image_text = to_markdown_image_text(upload_result.file_path, cleaned_alt)
                        markdown.append(f"{image_text}\n\n")
                        images.append(upload_result.file_path)
                        logger.info(f"Successfully saved image {upload_result.file_path}.")
                except Exception as ex:
                    logger.exception(f"Failed to save image: {upload_result.file_path}. Error: {ex}")
        return markdown, images

    def _extract_shape(self, slide_number, shape, save_name_template: str, tenant_id: str):
        markdown = []
        images = []
        if shape.name.startswith("Title"):
            # 标题
            markdown.append(f"# {shape.text}\n\n")
        elif shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART):
            # 图片
            new_markdown, new_images = self._extract_image_from_shape(shape, save_name_template, tenant_id)
            markdown.extend(new_markdown)
            images.extend(new_images)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # 文本框
            markdown.append(f"{shape.text}\n\n")
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # 表格
            table = shape.table
            markdown.append(self._convert_table_to_pai_table(table))
            markdown.append("\n\n")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                md, new_images = self._extract_shape(
                    slide_number, p, save_name_template, tenant_id
                )
                if md:
                    texts.append(md)
                    images.extend(new_images)
            markdown.append("\n".join(texts))
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            placeholder_type = shape.placeholder_format.type
            text = shape.text.strip()

            if placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.VERTICAL_TITLE,PP_PLACEHOLDER.HEADER):
                # 标题
                if text:
                    markdown.append(f"# {text}\n\n")
            elif placeholder_type == PP_PLACEHOLDER.BODY:
                if text:
                    markdown.append(f"{text}\n\n")
            elif placeholder_type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.CHART):
                # 图片
                new_markdown, new_images = self._extract_image_from_shape(shape, save_name_template, tenant_id)
                markdown.extend(new_markdown)
                images.extend(new_images)
            elif placeholder_type == PP_PLACEHOLDER.TABLE:
                if hasattr(shape, 'table') and shape.table:
                    table = shape.table
                    # 检查是否有至少 1 行 1 列，且内容非空
                    if len(table.rows) > 0 and len(table.columns) > 0:
                        # 可选：进一步检查是否有非空单元格
                        has_content = any(
                            cell.text.strip() 
                            for row in table.rows 
                            for cell in row.cells
                        )
                        if has_content:
                            markdown.append(self._convert_table_to_pai_table(table))
                            markdown.append("\n\n")


        return "".join(markdown), images

    def _convert_table_to_pai_table(self, table):
        table_matrix = [
            ["" for _ in range(len(table.columns))] for _ in range(len(table.rows))
        ]
        visited_cells = set()
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
                if (i, j) in visited_cells:
                    continue
                cell_content = table.cell(i, j).text.replace("\n", "").replace("\r", "")
                if table.cell(i, j).is_merge_origin:
                    col_span = table.cell(i, j).span_width
                    row_span = table.cell(i, j).span_height
                    while (
                        col_span > 1
                        and j + col_span <= len(table.columns)
                        and table_matrix[i][j + col_span - 1] == ""
                    ):
                        col_span -= 1
                        table_matrix[i][j + col_span] = cell_content
                        visited_cells.add((i, j + col_span))
                    while (
                        row_span > 1
                        and i + row_span <= len(table.rows)
                        and table_matrix[i + row_span - 1][j] == ""
                    ):
                        row_span -= 1
                        table_matrix[i + row_span][j] = cell_content
                        visited_cells.add((i + row_span, j))
                if table_matrix[i][j] == "":
                    table_matrix[i][j] = cell_content
                    visited_cells.add((i, j))

        row_headers_index = []
        col_headers_index = []
        if table.first_row:
            row_headers_index.append(0)
        if table.first_col:
            col_headers_index.append(0)
        pai_table = PaiTable(
            data=table_matrix,
            row_headers_index=row_headers_index,
            column_headers_index=col_headers_index,
        )
        return convert_table_to_markdown(pai_table, len(table.columns))

    def convert_pptx_to_markdown(
        self, presentation: Presentation, save_name_template: str, tenant_id: str
    ):
        markdown = []
        images = []
        slide_image_flag = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            image_flag = False
            for shape in slide.shapes:
                shape_markdown, shape_images = self._extract_shape(
                    slide_number, shape, save_name_template, tenant_id
                )
                markdown.append(shape_markdown)
                images.extend(shape_images)
            markdown.append(f"# slide_number_{slide_number}\n\n")
            slide_image_flag.append(image_flag)

        return "".join(markdown), images

    def read(self, file_item: FileItem) -> List[Document]:
        """
        Read a CSV file and return a list of Documents.
        """
        try:
            file_item.file.seek(0)
            presentation = Presentation(file_item.file)

            markdown_content, images = self.convert_pptx_to_markdown(
                presentation, file_item.kb_id + "/images/{}", tenant_id=file_item.tenant_id
            )

            metadata = file_item.metadata()
            markdown_content = replace_consecutive_spaces(markdown_content)
            docs = [Document(id_=file_item.id, text=markdown_content, metadata=metadata)]
            logger.info(f"Successfully read {file_item.file_name}.")

            return docs
        except Exception as e:
            logger.exception(e)
            return []
